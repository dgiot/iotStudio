# -*- coding: utf-8 -*-
"""dgiot_lite 时序数据采集管理 — 会议演示 PPT 生成器"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = r"D:\ai\dgiot_lite\docs\报价\对外\pdf\商务文档\时序数据采集管理系统-演示.pptx"
BLUE = RGBColor(0x0F, 0x23, 0x47)
CYAN = RGBColor(0x00, 0xB4, 0xD8)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tf

def add_line(slide, y, color=GOLD, width=10):
    line = slide.shapes.add_shape(1, Inches(1.67), Inches(y), Inches(width), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()
    return line

# === Slide 1: 封面 ===
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
add_text(s1, 2, 1.8, 9, 1.5, "时序数据采集与应用管理系统", 40, WHITE, True, PP_ALIGN.CENTER)
add_line(s1, 3.4, GOLD, 10)
add_text(s1, 2, 3.8, 9, 0.8, "Time-series Data Collection & Application Management", 20, CYAN, False, PP_ALIGN.CENTER)
add_line(s1, 4.7, GOLD, 10)
add_text(s1, 2, 5.2, 9, 0.6, "演示汇报  ·  2026.08", 16, GRAY, False, PP_ALIGN.CENTER)

# === Slide 2: 系统概述 ===
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, BLUE)
add_text(s2, 0.5, 0.3, 12, 0.6, "系统概述", 28, WHITE, True)
add_line(s2, 1.0, GOLD, 12)
items = [
    ("定位", "时序数据采集与应用管理 — 面向大庆油田 16 厂 · 100+ 作业区 · 260 万+ 点位"),
    ("架构", "中心侧 9 模块 + 边缘中枢 6 引擎 + 边缘采集 5 能力 · 三层部署"),
    ("协议", "10+ 工业协议 (A11/Modbus/OPC DA UA/IEC104/MQTT/HTTP/DTU/RTSP)"),
    ("算法", "15 种流式计算算法 · QPS 156 · 四级告警 · 定制注册"),
    ("交付", "29 前端页面 · 70 后端 API · 20+ PDF 文档 · 一键部署")
]
for i, (k, v) in enumerate(items):
    y = 1.5 + i * 1.1
    add_text(s2, 0.8, y, 2, 0.5, k, 18, GOLD, True)
    add_text(s2, 3.0, y, 9.5, 0.5, v, 16, WHITE)

# === Slide 3: 六域对齐 ===
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, BLUE)
add_text(s3, 0.5, 0.3, 12, 0.6, "六域功能架构 (申报 48 计价项对齐)", 28, WHITE, True)
add_line(s3, 1.0, GOLD, 12)
domains = [
    ("① 采得全", "10 协议 · 双路采集 · 10 家网关 · 断网补传", "100%"),
    ("② 存得住", "TDengine+PG+SQLite 三级存储 · 自动归档 · 8.3x 压缩", "100%"),
    ("③ 算得准", "15 种算法 · 四级告警 · 流式引擎 QPS 156", "100%"),
    ("④ 看得见", "29 页面 · GIS 地图 · 数字大屏 · 中英日俄西阿 7 语", "100%"),
    ("⑤ 管得住", "国密 SM2/3/4 · RBAC 3 角色 · 审计 3847 条 · 信创 5 平台", "100%"),
    ("⑥ 交付保障", "压测 360h · 灌数 100 万点 · 一键部署 · 四层校验", "100%"),
]
for i, (name, desc, score) in enumerate(domains):
    y = 1.4 + i * 0.95
    add_text(s3, 0.8, y, 3, 0.5, name, 18, GOLD, True)
    add_text(s3, 3.5, y, 7, 0.5, desc, 15, WHITE)
    add_text(s3, 11.5, y, 1.5, 0.5, score, 18, CYAN, True)

# === Slide 4: 核心功能 ===
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, BLUE)
add_text(s4, 0.5, 0.3, 12, 0.6, "核心功能矩阵", 28, WHITE, True)
add_line(s4, 1.0, GOLD, 12)
features = [
    ("数字大屏", "KPI卡片 · 拓扑总览 · 告警滚动 · 秒级刷新", "29 页面"),
    ("设备管理", "566 设备台账 · 六页签详情 · 五态生命周期", "设备组"),
    ("通道管理", "21 通道 · Modbus 地址段盲扫 · CRC 报文跟踪", "设备组"),
    ("流式计算", "15 算法 · 实时自增 · 6 条数据流 · QPS 156", "计算组"),
    ("影子设备", "Desired/Reported/Delta · 对标 dgaiot OTP", "设备详情"),
    ("国密安全", "SM2/3/4 · RBAC 3 角色 · 信创 5 平台验证", "系统组"),
    ("FDE 向导", "物模型→本体编译→扫描→规则→驾驶舱→AI 六步", "计算组"),
    ("国际化", "中/英/日/俄/西/阿 7 语 · 顶栏一键切换", "全局"),
]
for i, (name, desc, group) in enumerate(features):
    col = i % 4; row = i // 4
    x = 0.5 + col * 3.1; y = 1.4 + row * 1.4
    add_text(s4, x, y, 2.8, 0.4, name, 16, GOLD, True)
    add_text(s4, x, y + 0.4, 2.8, 0.6, desc, 12, WHITE)
    add_text(s4, x, y + 0.9, 2.8, 0.3, group, 10, GRAY)

# === Slide 5: 关键数字 ===
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, BLUE)
add_text(s5, 0.5, 0.3, 12, 0.6, "关键数字", 28, WHITE, True)
add_line(s5, 1.0, GOLD, 12)
nums = [("566", "设备在线"), ("21", "通道"), ("10", "协议"), ("15", "算法"), ("70", "API"), ("29", "页面"), ("7", "语言"), ("99.96%", "成功率")]
for i, (n, label) in enumerate(nums):
    x = 0.5 + (i % 4) * 3.1; y = 1.5 + (i // 4) * 2.5
    add_text(s5, x, y, 2.8, 1, n, 52, CYAN, True, PP_ALIGN.CENTER)
    add_text(s5, x, y + 1.2, 2.8, 0.5, label, 18, WHITE, False, PP_ALIGN.CENTER)

# === Slide 6: 部署拓扑 ===
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, BLUE)
add_text(s6, 0.5, 0.3, 12, 0.6, "部署拓扑 — 三层架构", 28, WHITE, True)
add_line(s6, 1.0, GOLD, 12)
layers = [
    ("中心侧 · 办公网", "数字大屏 · 设备管理 · 产品管理 · 通道管理 · 边缘代理\n设备模拟 · 流式计算 · 规则编排 · 运维任务 · 数据报表\n时序分析 · 链路拓扑 · 采集场景 · 用户管理 · GIS 地图"),
    ("边缘中枢 · DMZ 区 (麒麟 Linux)", "EMQX Broker · Parse Server · TDengine 时序库\nPostgreSQL 关系库 · NestJS 业务 · Vite 前端"),
    ("边缘采集 · 生产网 IO 服务器", "10+ 协议解析引擎 · 双路采集 · 兼容 A11\nModbus 盲扫 · 物模型自动配点 · 断网补传"),
]
for i, (name, desc) in enumerate(layers):
    y = 2 + i * 1.8
    add_text(s6, 0.8, y, 4, 0.5, name, 18, GOLD, True)
    add_text(s6, 5, y, 7.5, 1.5, desc, 14, WHITE)
    if i < 2:
        add_text(s6, 2.5, y + 1.2, 1, 0.5, "↓", 24, CYAN, False, PP_ALIGN.CENTER)

# === Slide 7: 技术栈 ===
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, BLUE)
add_text(s7, 0.5, 0.3, 12, 0.6, "技术栈", 28, WHITE, True)
add_line(s7, 1.0, GOLD, 12)
tech = [("后端", "Python FastAPI + uvicorn · 70 API"), ("前端", "Vue 3 + Vite + Element Plus + ECharts"),
        ("协议", "pymodbus · asyncua · 自研 A11/CommBridge"), ("存储", "TDengine + PostgreSQL + SQLite"),
        ("消息", "MQTT (dgiot/EMQX Broker)"), ("安全", "国密 SM2/SM3/SM4 · TLS 1.2"),
        ("部署", "Docker · Nginx · Windows/Linux"), ("国际化", "vue-i18n · 7 语言")]
for i, (k, v) in enumerate(tech):
    col = i % 2; row = i // 2
    x = 0.5 + col * 6.2; y = 1.5 + row * 1.3
    add_text(s7, x, y, 1.5, 0.4, k, 18, GOLD, True)
    add_text(s7, x + 1.7, y, 4.5, 0.4, v, 15, WHITE)

# === Slide 8: Demo ===
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8, BLUE)
add_text(s8, 0.5, 0.3, 12, 0.6, "Demo 演示", 28, WHITE, True)
add_line(s8, 1.0, GOLD, 12)
add_text(s8, 1, 2.5, 11, 0.8, "http://localhost:20036", 36, CYAN, True, PP_ALIGN.CENTER)
add_text(s8, 1, 3.5, 11, 0.5, "账号: dgiot_dev  /  密码: dgiot_dev", 20, WHITE, False, PP_ALIGN.CENTER)
add_text(s8, 1, 4.5, 11, 0.5, "云端: dev.dgiotcloud.cn:5180", 18, GRAY, False, PP_ALIGN.CENTER)
add_text(s8, 1, 5.5, 11, 0.5, "六域 100% 对齐  ·  70 API  ·  29 页面  ·  7 语言", 16, GOLD, False, PP_ALIGN.CENTER)

# === Slide 9: 谢谢 ===
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9)
add_text(s9, 2, 2.5, 9, 1.5, "谢谢", 52, GOLD, True, PP_ALIGN.CENTER)
add_text(s9, 2, 4, 9, 0.8, "时序数据采集与应用管理系统", 24, WHITE, False, PP_ALIGN.CENTER)
add_text(s9, 2, 5, 9, 0.5, "2026.08", 16, GRAY, False, PP_ALIGN.CENTER)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
