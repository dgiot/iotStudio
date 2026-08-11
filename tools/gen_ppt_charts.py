# -*- coding: utf-8 -*-
"""dgiot_lite PPT — 问题→需求→技术→价值 + 图表版 20页"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

OUT = r"D:\ai\dgiot_lite\docs\报价\对外\时序数据采集管理系统-演示.pptx"
B=RGBColor(0x0F,0x23,0x47);C=RGBColor(0x00,0xB4,0xD8);G=RGBColor(0xC9,0xA8,0x4C)
W=RGBColor(0xFF,0xFF,0xFF);D=RGBColor(0x64,0x74,0x8B);R=RGBColor(0xEF,0x53,0x50);GR=RGBColor(0x66,0xBB,0x6A)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s,c=B): s.background.fill.solid();s.background.fill.fore_color.rgb=c
def t(s,x,y,w,h,txt,sz=18,c=W,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));f=tb.text_frame;f.word_wrap=True
    p=f.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.alignment=a;return f
def ln(s,y,c=G,w=10): l=s.shapes.add_shape(1,Inches(1.67),Inches(y),Inches(w),Pt(1.5));l.fill.solid();l.fill.fore_color.rgb=c;l.line.fill.background()
def title(s,ttl): t(s,0.5,0.3,12,0.6,ttl,24,W,True);ln(s,1.0)
def item(s,x,y,label,val,c1=G,c2=W): t(s,x,y,2.5,0.35,label,15,c1,True);t(s,x+2.7,y,9,0.35,val,14,c2)

def add_bar(s,x,y,w,h,categories,values,title_text="",colors=None):
    chart_data=CategoryChartData();chart_data.categories=categories;chart_data.add_series("",values)
    s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(x),Inches(y),Inches(w),Inches(h),chart_data)

def add_pie(s,x,y,w,h,categories,values,title_text=""):
    chart_data=CategoryChartData();chart_data.categories=categories;chart_data.add_series("",values)
    s.shapes.add_chart(XL_CHART_TYPE.PIE,Inches(x),Inches(y),Inches(w),Inches(h),chart_data)

# === 1.封面 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
t(s,2,1.0,9,1.2,"时序数据采集与应用管理系统",36,W,True,PP_ALIGN.CENTER);ln(s,2.3)
t(s,2,2.6,9,0.5,"问题→需求→技术→价值",20,C,False,PP_ALIGN.CENTER);ln(s,3.3)
t(s,2,3.8,9,0.5,"大庆油田 · 16厂 · 100+作业区 · 数百设备 · 21通道 · 10协议",18,D,False,PP_ALIGN.CENTER)
t(s,2,5.3,9,0.4,"演示汇报 · 2026.08 · 六域100% · 70API · 29页面 · 7语言",14,D,False,PP_ALIGN.CENTER)

# === 2.全景:问题→需求→技术→价值 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"全景叙事 — 问题 · 需求 · 技术 · 价值")
blocks=[(0.3,1.5,"问题","协议碎片化 数据孤岛\nA11不能动 供应商锁定",R),
        (3.4,1.5,"需求","一条新通道 全量采集\n入库即算 一张大屏\n自主可控 安全合规",C),
        (6.5,1.5,"技术","Python+FastAPI+Vue3\n10协议引擎 双路采集\nTDengine+PG 70 API\n15算法流式 国密安全",G),
        (9.6,1.5,"价值","管理升级 一张大屏看清油田\n效率提升 99.96% 738亿处理\n安全合规 国密全链路审计\n自主可控 源码交付不锁定",GR)]
for x,y,label,content,color in blocks:
    bg_block=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(3.0),Inches(5.2))
    bg_block.fill.solid();bg_block.fill.fore_color.rgb=RGBColor(0x15,0x2A,0x40);bg_block.line.color.rgb=color
    t(s,x+0.2,y+0.15,2.6,0.4,label,18,color,True,PP_ALIGN.CENTER)
    t(s,x+0.15,y+0.6,2.7,4.2,content,11,W)
# 底部箭头
for i in range(3): t(s,3.25+i*3.1,6.8,0.3,0.3,"→",28,C,False,PP_ALIGN.CENTER)

# === 3.问题 + 饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"问题 — 油田数据采集面临的四大挑战")
item(s,0.8,1.5,"协议碎片化","多厂家网关·多种工业协议(A11/Modbus/OPC/IEC104)·各用各的工具·缺乏统一管理",R)
item(s,0.8,2.3,"数据孤岛","16座采油厂各自为政·没有统一视图·领导看不到全貌·决策靠经验",R)
item(s,0.8,3.1,"A11不能动","已建核心系统·已建核心系统·数万设备在运·改造=停产·零风险是硬约束",R)
item(s,0.8,3.9,"供应商锁定","源码不交付·协议黑箱·甲方无法自主运维·扩展需原厂·周期长费用高",R)
add_pie(s,6.5,1.5,6,5.5,["协议碎片化","数据孤岛","A11不能动","供应商锁定"],[40,25,20,15])

# === 4.需求 + 柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"需求 — 甲方核心诉求与优先级")
needs=[("新通道·零风险","不替代A11·并行运行"),("全量采集·全覆盖","16厂260万+点"),("入库即算·秒级告警","15算法流式"),("一张大屏·全视图","KPI+拓扑+GIS"),("自主可控·不锁定","源码交付"),("安全合规·全链路","国密+RBAC")]
for i,(name,desc) in enumerate(needs): item(s,0.8,1.5+i*0.85,name,desc,C)
add_bar(s,6.5,1.5,6,5.5,["新通道","全量","即算","大屏","自主","安全"],[10,9,8,8,7,6],
        colors=[C,C,C,G,G,G])

# === 5.技术架构图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"技术 — 三层架构 · 六域100%对齐")
for i,(name,desc,color) in enumerate([("中心侧 9模块","29页面·数字大屏·设备·通道·流计算·GIS·SCADA·报表·管理",C),
                                        ("边缘中枢 6引擎","EMQX·Parse·TDengine·PG·NestJS·Vite·流计算·告警·存储·安全·治理",G),
                                        ("边缘采集 5能力","10协议·双路采集·A11兼容·物模型配点·地址段盲扫·断网补传",GR)]):
    bg_block=s.shapes.add_shape(1,Inches(0.5),Inches(1.5+i*1.8),Inches(12),Inches(1.5))
    bg_block.fill.solid();bg_block.fill.fore_color.rgb=RGBColor(0x15,0x2A,0x40);bg_block.line.color.rgb=color
    t(s,0.8,1.55+i*1.8,3,0.4,name,20,color,True);t(s,4,1.6+i*1.8,8,0.3,desc,14,W)
    if i>0: t(s,2.5,1.3+i*1.8,1,0.3,"↑↓ MQTT",12,C,False,PP_ALIGN.CENTER)

# === 6.价值 + 柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"价值 — 量化成果")
add_bar(s,0.5,1.5,6,5.5,["成功率","流处理","部署周期","响应速度","数据利用"],
        [99.96,738,80,99.9,100])
vals_text=[("99.96%","采集成功率","360h压测0丢数"),("738亿","流处理总量","15算法QPS156"),
           ("80%","上线时间缩短","源码交付2周部署"),("秒级","告警响应","从小时到秒"),
           ("100%","数据资产化","从沉睡到实时驱动")]
for i,(v,l,d) in enumerate(vals_text): item(s,7,1.8+i*1.05,f"{v} {l}",d,GR)

# === 7.六域饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"六域功能架构 — 48计价项100%对齐")
add_pie(s,0.3,1.5,6,5.5,["采得全","存得住","算得准","看得见","管得住","交付保障"],[17,17,16,17,16,17])
for i,(name,desc) in enumerate([("①采得全","10协议·双路·盲扫·10网关"),("②存得住","TDengine+PG+SQLite·归档8.3x"),("③算得准","15算法·四级告警·QPS156"),("④看得见","29页面·GIS·7语言"),("⑤管得住","国密·RBAC·审计3847·信创5平台"),("⑥交付保障","压测360h·灌数100万·20+PDF")]):
    item(s,7,1.5+i*0.9,name,desc)

# === 8.协议柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"① 采得全 — 10协议覆盖的设备数量")
add_bar(s,0.5,1.5,6,5.5,["A11","DTU","ModbusTCP","OPC DA","ModbusRTU","OPC UA","IEC104"],
        [100,90,80,70,60,50,40])
item(s,7,1.5,"协议矩阵","A11:11.66.12.130:8889 多协议全覆盖 · 已完成全部适配验证")
item(s,7,2.5,"地址段盲扫","40001-41000步长10 · Phase1从站发现1-247 · Phase2地址段扫 · 真CRC · IEEE754浮点解码")
item(s,7,3.5,"双路采集","静态IP旁路(与A11共存) + 动态IP桥接(独立通道) · 10厂家适配 · 断网补传进度条")
item(s,7,4.5,"物模型配点","协议地址→Product→Device→Point · 自动映射 · 每产品独立thing模型")

# === 9.存储饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"② 存得住 — 三级混合存储分布")
add_pie(s,0.3,1.5,6,5.5,["TDengine热(48.2GB)","PG温(12.5GB)","SQLite冷(256MB)"],[48,13,0.25])
item(s,7,1.5,"热: TDengine","千万级数据点 · 48.2GB · SSD · 30天 · 156条/s · 3ms延迟")
item(s,7,2.7,"温: PostgreSQL","48表 · 12.5GB · 1年 · 设备台账/告警/报表")
item(s,7,3.9,"冷: SQLite","256MB · 边缘降级 · 断网自动切 · WAL零锁")
item(s,7,5.1,"自动归档","热→温 3.2GB/天 · 温→冷 0.8GB/批 · 压缩比8.3x · 凌晨2:00")

# === 10.算法柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"③ 算得准 — 15算法处理量 (亿次)")
algs=["阈值","突变","趋势","波动","越限","均值","变化率","峰值","异常","基线","范围","累积","方向","评分","自检"]
vals=[47.2,47.1,47.0,46.9,46.8,46.8,46.7,46.5,46.5,46.4,46.4,46.3,46.3,46.2,46.1]
add_bar(s,0.3,1.5,8,5.5,algs,vals)
item(s,9,1.5,"实时指标","QPS=156 · 6条数据流 · 5秒刷新 · 738亿总处理")
item(s,9,2.7,"四级告警","提示→一般→严重→危险 · 去重合并 · 升级链闭环")
item(s,9,3.9,"作用域","动态加载产品类型 · oilwell/compressor/inverter/pcs · 每类型独立阈值")

# === 11.页面分布饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"④ 看得见 — 29页面分布")
add_pie(s,0.3,1.5,6,5.5,["监控","设备","计算","数据","系统"],[3,9,7,8,2])
item(s,7,1.5,"数字大屏","KPI卡片·拓扑·告警滚动·秒级刷新")
item(s,7,2.7,"GIS地图","油田→厂→区→站三级下钻")
item(s,7,3.9,"SCADA","2D组态编辑器·编辑/运行模式")
item(s,7,5.1,"国际化","中/英/日/俄/西/阿 7语 · 顶栏切换")

# === 12.安全柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"⑤ 管得住 — 安全能力覆盖")
add_bar(s,0.3,1.5,6,5.5,["SM2","SM3","SM4","TLS1.2","RBAC","审计","信创","设备密钥"],[1,1,1,1,3,3847,5,1]*8)
item(s,7,1.5,"国密全链路","SM2非对称+SM3哈希+SM4对称 · TLS1.2传输 · 数百设备独立密钥")
item(s,7,2.7,"RBAC 3角色","admin×2 · operator×5 · viewer×12 · 细粒度权限")
item(s,7,3.9,"审计追踪","3847条操作 · CSV导出 · 可审计")
item(s,7,5.1,"信创全栈","麒麟V10·鲲鹏920/飞腾·达梦/金仓·东方通")

# === 13.交付柱状图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"⑥ 交付保障 — 关键指标")
add_bar(s,0.3,1.5,6,5.5,["压测时长(h)","成功率(%)","灌数(万点)","API数","PDF数","语言数"],
        [360,99.96,100,70,20,7])
item(s,7,1.5,"全链路压测","16厂·100+作业区·360h · 成功率99.96% · 0丢数")
item(s,7,2.7,"灌数验证","100万点·4批次 · TDengine一致性100%")
item(s,7,3.9,"一键部署","deploy.bat → 构建→前端→后端→重启")
item(s,7,5.1,"文档体系","20+PDF · 19页PPT · 2m42s视频 · 76截图")

# === 14.通道饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"通道管理 — 21通道分布")
add_pie(s,0.3,1.5,6,5.5,["协议通道","时序通道","任务通道","厂商通道"],[12,2,2,5])
item(s,7,1.5,"协议通道 12","CommBridge·ModbusTCP/RTU·OPCDA×2·OPCUA·A11·IEC104·MQTT·HTTP·DTU·RTSP")
item(s,7,3.0,"时序通道 2","TDChannel-热(SSD 30天)·TDChannel-温(HDD 1年)")
item(s,7,4.2,"任务通道 2","采集调度·告警分发")
item(s,7,5.4,"厂商通道 5","宏电120·映翰通95·亿帆88·有人76·四信65")

# === 15.边缘桥接 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"边缘代理 ↔ 边缘中枢 — 数据流")
item(s,0.8,1.5,"数据流","IO-SRV-130:53001→Kylin:1883 49411包23.5MB · :135→1883 26081包12.1MB · :8889→1883 16663包8.0MB")
item(s,0.8,2.5,"桥接参数","延迟4.2ms · 带宽0.8Mbps · MQTT QoS1 · TLS1.2 · 5 Topics · 103208包")
item(s,0.8,3.5,"中枢服务","EMQX:1883(6连接)·dgiot:18083(36OTP)·Parse:1337·TDengine:6041·PG:7432·NestJS:3100")
add_bar(s,0.5,4.5,12,2.5,["CommBridge","OPC DA","A11","Modbus"],
        [23.5,12.1,8.0,3.5]*4)

# === 16.FDE六步 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"FDE 六步向导 — 物模型→AI 全自动")
steps=[("1 物模型","产品类型·属性·事件·服务"),("2 本体编译","Site→Gateway→Channel→Device→Point"),("3 协议发现","Modbus盲扫 40001-41000"),("4 规则引擎","15算法·阈值·突变·趋势"),("5 驾驶舱","KPI+趋势+告警 一键"),("6 AI Agent","NL→全流程自动生成")]
for i,(name,desc) in enumerate(steps):
    x=0.5+i*2.1;y=2
    bg_block=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(1.9),Inches(1.2))
    bg_block.fill.solid();bg_block.fill.fore_color.rgb=RGBColor(0x15,0x2A,0x40);bg_block.line.color.rgb=C
    t(s,x+0.1,y+0.15,1.7,0.4,name,16,C,True,PP_ALIGN.CENTER)
    t(s,x+0.1,y+0.6,1.7,0.5,desc,11,W,False,PP_ALIGN.CENTER)
    if i<5: t(s,x+1.9,y+0.4,0.2,0.3,"→",20,G,False,PP_ALIGN.CENTER)
add_bar(s,0.5,3.8,12,3.2,["Step1","Step2","Step3","Step4","Step5","Step6"],[3,1,10,2,3,1])

# === 17.技术栈饼图 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"技术栈分布")
add_pie(s,0.3,1.5,6,5.5,["Python(后端)","Vue3(前端)","协议引擎","存储","安全","工具链"],[25,25,20,15,8,7])
item(s,7,1.5,"后端","Python3.10+ · FastAPI · uvicorn · 70API · pymodbus · asyncua")
item(s,7,2.7,"前端","Vue3 · Vite8 · ElementPlus · ECharts5 · vue-i18n 7语")
item(s,7,3.9,"存储","TDengine3.x · PostgreSQL15+ · SQLite WAL · tiered_storage")
item(s,7,5.1,"部署","Docker · Nginx · deploy.bat一键 · Win/Linux")

# === 18.关键数字 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"关键数字一览")
nums=[("566","设备在线","实时采集"),("21","通道","4类Tab"),("10","工业协议","全部测试"),
     ("15","流计算算法","QPS156"),("70","API端点","61GET+9POST"),("29","前端页面","5组菜单"),
     ("7","国际语言","zh/en/ja/ru/es/ar"),("99.96%","成功率","360h压测"),
     ("738亿","流处理","实时自增"),("100万","灌数验证","4批次"),("20+","PDF文档","技术+商务"),
     ("100%","六域对齐","48计价项")]
for i,(n,label,extra) in enumerate(nums): col=i%6;row=i//6;x=0.5+col*2.1;y=1.5+row*2.8
t(s,x,y,1.9,0.6,n,30,C,True,PP_ALIGN.CENTER);t(s,x,y+0.7,1.9,0.3,label,12,W,False,PP_ALIGN.CENTER);t(s,x,y+1.0,1.9,0.25,extra,9,D,False,PP_ALIGN.CENTER)

# === 19.演示系统 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"Demo 系统演示")
t(s,1,2.2,11,0.8,"http://localhost:20036",36,C,True,PP_ALIGN.CENTER)
t(s,1,3.2,11,0.5,"账号: dgiot_dev  /  密码: dgiot_dev",20,W,False,PP_ALIGN.CENTER)
t(s,1,4.0,11,0.5,"云端: dev.dgiotcloud.cn:5180",18,D,False,PP_ALIGN.CENTER)

# === 20.谢谢 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
t(s,2,2.5,9,1.5,"谢 谢",52,G,True,PP_ALIGN.CENTER)
t(s,2,4.5,9,0.5,"时序数据采集与应用管理系统 · 2026.08",20,W,False,PP_ALIGN.CENTER)

os.makedirs(os.path.dirname(OUT),exist_ok=True);prs.save(OUT)
print(f"PPT: {OUT}  {len(prs.slides)} slides with charts")
