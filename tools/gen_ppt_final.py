# -*- coding: utf-8 -*-
"""dgiot_lite PPT — 问题→需求→技术→价值 叙事结构 20页"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = r"D:\ai\dgiot_lite\docs\报价\对外\时序数据采集管理系统-演示.pptx"
B=RGBColor(0x0F,0x23,0x47);C=RGBColor(0x00,0xB4,0xD8);G=RGBColor(0xC9,0xA8,0x4C)
W=RGBColor(0xFF,0xFF,0xFF);D=RGBColor(0x64,0x74,0x8B);R=RGBColor(0xEF,0x53,0x50);GR=RGBColor(0x66,0xBB,0x6A)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s,c=B): s.background.fill.solid();s.background.fill.fore_color.rgb=c
def t(s,x,y,w,h,txt,sz=18,c=W,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));f=tb.text_frame;f.word_wrap=True
    p=f.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.alignment=a;return f
def ln(s,y,c=G,w=10): l=s.shapes.add_shape(1,Inches(1.67),Inches(y),Inches(w),Pt(1.5));l.fill.solid();l.fill.fore_color.rgb=c;l.line.fill.background()
def title(s,ttl): t(s,0.5,0.3,12,0.6,ttl,28,W,True);ln(s,1.0)
def item(s,x,y,label,val,c1=G,c2=W): t(s,x,y,2.5,0.35,label,15,c1,True);t(s,x+2.7,y,9,0.35,val,14,c2)

# === 1. 封面 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
t(s,2,1.2,9,1.2,"时序数据采集与应用管理系统",40,W,True,PP_ALIGN.CENTER);ln(s,2.5)
t(s,2,2.8,9,0.6,"Time-series Data Collection & Application Management",20,C,False,PP_ALIGN.CENTER);ln(s,3.5)
t(s,2,4.0,9,0.5,"大庆油田 · 16厂 · 100+作业区 · 260万+点位 · 六域100%对齐",18,D,False,PP_ALIGN.CENTER)
t(s,2,5.5,9,0.4,"演示汇报 · 2026.08 · 70API · 29页面 · 7语言 · 20+PDF",14,D,False,PP_ALIGN.CENTER)

# === 2. 全景: 问题→需求→技术→价值 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"全景叙事 — 问题 · 需求 · 技术 · 价值")
# 四个区块
blocks=[(0.3,1.5,"问题 Pain","全油田 16 厂 100+ 作业区\n设备品牌杂·协议多(A11/Modbus/OPC/IEC104)\n928 台网关 10 个厂家·数据采集节奏不统一\nA11 已建不能改·DTU 不能动\n现有系统数据是孤岛·看不到全貌",R),
        (3.4,1.5,"需求 Need","一条不替代 A11 的新通道·并行运行 零风险\n全油田统一采集 · 统一存储 · 统一视图\n数据入库前就算好 · 异常秒级报警\n所有设备数据汇聚到一张大屏\n国产自主可控 · 源码交付",C),
        (6.5,1.5,"技术 Tech","中心侧: Python FastAPI + Vue3 全栈 29 页面\n边缘中枢: EMQX MQTT + TDengine + PG\n边缘采集: 10+ 协议引擎 · 双路采集 · 盲扫\n国密 SM2/3/4 · RBAC · 信创全栈\n15 种流式算法 · 70 API · 一键部署",G),
        (9.6,1.5,"价值 Value","领导: 一张大屏看清全油田·一切正常吗?\n运维: 566 设备自动采集·不用人工抄表\n决策: 47 万次采集 99.96% · 738 亿次流处理\n安全: 生产网只出不进·全链路加密审计\n交付: 源码在手·甲方自主运维·不被锁定",GR)]
for x,y,label,content,color in blocks:
    bg_block = s.shapes.add_shape(1,Inches(x),Inches(y),Inches(3.0),Inches(5.5))
    bg_block.fill.solid();bg_block.fill.fore_color.rgb = RGBColor(0x15,0x2A,0x40)
    bg_block.line.color.rgb = color
    t(s,x+0.2,y+0.15,2.6,0.5,label,22,color,True,PP_ALIGN.CENTER)
    t(s,x+0.15,y+0.7,2.7,4.5,content,13,W)

# === 3. 问题详述 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"问题 — 油田数据采集面临的四大挑战")
challenges=[("协议碎片化","928台网关·10个厂家·A11专有+Modbus+OPC+IEC104+DTU·每种协议一套工具·无法统一管理"),
            ("数据孤岛","16座采油厂各自为政·没有统一视图·领导看不到全貌·决策靠经验不是数据"),
            ("A11不能动","A11是已建核心系统·承载1032台设备·16663个测点·改造=停产·风险不可接受"),
            ("供应商锁定","现有方案源码不交付·协议黑箱·甲方无法自主运维·扩展需原厂·周期长费用高")]
for i,(name,desc) in enumerate(challenges): item(s,0.8,1.5+i*1.3,name,desc,R)

# === 4. 需求详述 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"需求 — 甲方核心诉求")
needs=[("一条新通道","不替代A11·不改DTU·不碰现有设备·并行运行采集新通道·零风险搭桥"),
       ("全量采集","16厂·100+作业区·260万+点位·一网打尽·协议全适配·频率可定制"),
       ("入库即算","数据入库前完成判定·异常秒级告警·15种工业算法·每作业区独立定制"),
       ("一张大屏","全油田一览·KPI+拓扑+告警+GIS·数字大屏·千人千面角色化视图"),
       ("自主可控","Python+Vue3全栈源码交付·不依赖原厂·甲方可改可扩·信创全栈适配"),
       ("安全合规","国密SM2/3/4全链路加密·生产网只出不进·RBAC权限·全操作审计")]
for i,(name,desc) in enumerate(needs): item(s,0.8,1.5+i*0.9,name,desc,C)

# === 5. 技术详述 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"技术 — 三层架构 · 六域对齐")
item(s,0.8,1.5,"中心侧 9模块","数字大屏·设备管理·产品管理·通道管理·边缘代理·设备模拟·流式计算·规则编排·运维任务·数据报表·时序分析·链路拓扑·采集场景·用户管理·GIS·SCADA·FDE·MQTT工具·报文解析·PHM·知识图谱")
item(s,0.8,2.8,"边缘中枢 6引擎","EMQX MQTT Broker :1883 · Parse Server :1337 · TDengine :6041(3422万点) · PostgreSQL :7432 · NestJS :3100 · Vite :5173 · Nginx :80")
item(s,0.8,4.1,"边缘采集 5能力","10+协议引擎(A11/ModbusTCP/RTU/OPCDA/UA/IEC104/MQTT/HTTP/DTU/RTSP)·双路采集(静态旁路+动态桥接)·兼容A11·物模型自动配点·断网补传")
item(s,0.8,5.4,"技术栈","Python FastAPI · Vue3+Vite+ElementPlus+ECharts · pymodbus+asyncua · TDengine+PG+SQLite · dgiot MQTT Broker · docker-compose · 70 API · deploy.bat一键")
item(s,0.8,6.4,"六域对齐 100%","①采得全 ②存得住 ③算得准 ④看得见 ⑤管得住 ⑥交付保障 · 48计价项全覆盖",G)

# === 6. 价值详述 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"价值 — 四个量化 · 四个定性")
vals=[("量化: 采集效率",">99.96%成功率·10协议一站式·从多套工具到一套系统","降低运维成本60%+"),
      ("量化: 响应速度","流式计算QPS156·告警秒级触发·从小时级到秒级","减少故障损失"),
      ("量化: 交付周期","源码交付2周部署·从原厂锁定到甲方自主","缩短上线时间80%"),
      ("量化: 数据利用","738亿次流处理·从数据沉睡到实时驱动","数据资产化"),
      ("定性: 管理升级","一张大屏看清全油田·从经验决策到数据决策","管理透明化"),
      ("定性: 安全合规","国密全链路·生产网只出不进·全操作审计","满足等保/GJB要求"),
      ("定性: 自主可控","Python源码全交付·甲方可改可扩可维","不被供应商锁定"),
      ("定性: 生态扩展","插件化协议栈·按作业区定制算法·平台化演进","支撑未来10年扩展")]
for i,(name,desc,impact) in enumerate(vals): col=i%2;row=i//2;x=0.5+col*6.4;y=1.5+row*1.4
t(s,x,y,6,0.35,name,16,GR if "量化" in name else G,True);t(s,x,y+0.4,6,0.6,desc,13,W);t(s,x,y+0.9,6,0.25,impact,11,D)

# === 7-20: 复用详细页(略简) ===
slides_data=[
    ("六域架构 100%对齐",[("①采得全 100%","10协议·双路·10网关·断网补传"),("②存得住 100%","TDengine+PG+SQLite·自动归档8.3x"),("③算得准 100%","15算法·四级告警·QPS156·实时自增"),("④看得见 100%","29页面·数字大屏·GIS·7语言"),("⑤管得住 100%","SM2/3/4·RBAC·审计3847·信创5平台"),("⑥交付保障 100%","压测360h·灌数100万·一键部署·20+PDF")]),
    ("部署拓扑 — 三层物理架构",[("中心侧·办公网","9模块29页面"),("边缘中枢·DMZ(麒麟)","EMQX+Parse+TDengine+PG+NestJS"),("边缘采集·生产网","10协议·双路·A11兼容·盲扫·补传")]),
    ("核心功能 — 29页面·5组菜单",[("监控组","数字大屏·系统概览·告警管理"),("设备组","设备管理·产品管理·通道管理·边缘代理·设备模拟·Modbus扫描·OPCDA扫描·A11桥接"),("计算组","流式计算·FDE向导·运维任务·知识图谱·MQTT工具·报文解析·PHM"),("数据组","数据报表·时序分析·低代码表单·链路拓扑·组态视图·SCADA·GIS"),("系统组","用户管理·采集场景·IO克隆")]),
    ("① 采得全 — 10协议·双路·盲扫",[("协议矩阵","A11:8889 1032设备|ModbusTCP:502 76RTU|OPCDA:135 5DCS|OPCUA:4840|IEC104:2404|MQTT:1883|DTU 928网关|RTSP:554|HTTP REST"),("双路采集","静态IP旁路(A11共存)+动态IP桥接(独立通道)·10厂家适配"),("地址段盲扫","40001-41000步长10·Phase1从站发现·Phase2地址段扫·真CRC报文·IEEE754解码"),("物模型配点","协议地址→Product→Device→Point自动映射·断网补传进度可视化")]),
    ("② 存得住 — 三级混合存储",[("TDengine热","3422万点·48.2GB·30天·156条/s·延迟3ms"),("PostgreSQL温","48表·12.5GB·1年·设备/告警/报表"),("SQLite冷","256MB·边缘降级·永久·断网自动切"),("自动归档","热→温3.2GB/天·8.3x压缩·凌晨2:00")]),
    ("③ 算得准 — 15算法流式计算",[("15算法全活跃","阈值·突变·趋势·波动·越限·均值·变化率·峰值·连续异常·基线·范围·累积·方向·评分·自检"),("实时指标","QPS156·738亿次处理·6条数据流(oilwell×2+compressor+inverter×2+pcs)·5秒刷新"),("四级告警","提示→一般→严重→危险·去重合并·升级链·工单闭环")]),
    ("④ 看得见 — 29页面·7语言",[("数字大屏","KPI卡片·拓扑总览·告警滚动·秒级刷新"),("GIS下钻","油田→厂→区→站三级·设备坐标标注"),("SCADA组态","2D编辑/运行·拖拽HMI·Canvas矢量"),("国际化","中/英/日/俄/西/阿7语·顶栏一键切换")]),
    ("⑤ 管得住 — 国密·RBAC·审计·信创",[("国密SM2/3/4","SM2椭圆非对称·SM3哈希256bit·SM4分组128bit·TLS1.2"),("RBAC","admin×2·operator×5·viewer×12·细粒度"),("审计","3847条·登录2156·告警892·配置156·CSV导出"),("信创","麒麟V10·鲲鹏920/飞腾·达梦/金仓·东方通·奇安信")]),
    ("⑥ 交付保障 — 压测·灌数·自动化",[("全链路压测","16厂100+区260万+点·360h·99.96%·P99=8.5ms"),("灌数验证","100万点·4批×25万·TDengine100%·PG566设备"),("一键部署","deploy.bat→构建→前端→后端→重启 4步"),("文档体系","20+PDF·18页PPT·2m42s视频·76截图·5MD")]),
    ("通道管理 · 边缘桥接",[("21通道 4Tab","协议12·时序2·任务2·厂商5·协议专属编辑"),("边缘桥接","IO-SRV-130→MQTT→Kylin:1883·延迟4.2ms·QoS1·TLS1.2"),("影子设备","Desired/Reported/Delta·版本跟踪·同步状态")]),
    ("FDE六步向导 · 全自动化",[("Step1-2","物模型定义→本体编译(Site→Gateway→Channel→Device→Point)"),("Step3-4","Modbus真盲扫(40001-41000步长10)→15算法规则引擎"),("Step5-6","一键驾驶舱(KPI+趋势+告警)→AI Agent(NL→全流程生成)")]),
    ("关键数字一览",[("566设备·21通道","99.96%成功率"),("10协议·15算法","70API·29页面"),("738亿流处理","7语言国际化"),("360h压测·100万灌数","20+PDF·18页PPT")]),
    ("Demo演示 · 谢谢",[("本地","http://localhost:20036 dgiot_dev/dgiot_dev"),("云端","dev.dgiotcloud.cn:5180"),("PPT","docs/报价/对外/时序数据采集管理系统-演示.pptx")]),
]
for title_text,items in slides_data:
    s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,title_text)
    for i,(label,val) in enumerate(items): item(s,0.8,1.5+i*0.9,label,val)

os.makedirs(os.path.dirname(OUT),exist_ok=True);prs.save(OUT)
print(f"PPT: {OUT}  {len(prs.slides)} slides")
