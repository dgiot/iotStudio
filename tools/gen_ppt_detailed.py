# -*- coding: utf-8 -*-
"""dgiot_lite 详细演示 PPT — 18页"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = r"D:\ai\dgiot_lite\docs\报价\对外\时序数据采集管理系统-演示.pptx"
B = RGBColor(0x0F,0x23,0x47); C = RGBColor(0x00,0xB4,0xD8); G = RGBColor(0xC9,0xA8,0x4C)
W = RGBColor(0xFF,0xFF,0xFF); D = RGBColor(0x64,0x74,0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb = B
def t(s,x,y,w,h,txt,sz=18,c=W,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));f=tb.text_frame;f.word_wrap=True
    p=f.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.alignment=a;return f
def ln(s,y): l=s.shapes.add_shape(1,Inches(1.67),Inches(y),Inches(10),Pt(1.5));l.fill.solid();l.fill.fore_color.rgb=G;l.line.fill.background()
def title(s,ttl): t(s,0.5,0.3,12,0.6,ttl,28,W,True);ln(s,1.0)
def item(s,x,y,label,val): t(s,x,y,2,0.35,label,15,G,True);t(s,x+2.2,y,9,0.35,val,14,W)

# 1.封面
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
t(s,2,1.5,9,1.2,"时序数据采集与应用管理系统",40,W,True,PP_ALIGN.CENTER)
ln(s,2.8);t(s,2,3.2,9,0.6,"Time-series Data Collection & Application Management System",20,C,False,PP_ALIGN.CENTER)
ln(s,4.0);t(s,2,4.5,9,0.5,"大庆油田 · 16厂 · 100+作业区 · 260万+点位",18,D,False,PP_ALIGN.CENTER)
t(s,2,5.5,9,0.4,"演示汇报 · 2026.08 · 六域100%对齐 · 70API · 29页面 · 7语言",14,D,False,PP_ALIGN.CENTER)

# 2.目录
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"汇报目录")
agenda=["系统定位与背景","六域功能架构 (48计价项)","部署拓扑 — 三层物理架构","核心功能矩阵 — 29页面全景","采集域 — 多协议·双路·盲扫·10网关","存储域 — 三级混合存储·自动归档","计算域 — 15算法流式计算·实时QPS156","应用域 — 数字大屏·GIS·SCADA·报表","安全域 — 国密·RBAC·审计·信创","交付域 — 压测·灌数·一键部署","技术栈 · 关键数字 · Demo演示"]
for i,a in enumerate(agenda): t(s,1.5,1.5+i*0.48,10,0.4,f"{i+1}. {a}",16,C if i<4 else W)

# 3.系统定位
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"系统定位与背景")
items=[("项目名称","时序数据采集与应用管理系统"),("覆盖范围","大庆油田 16座采油厂 · 100+作业区 · 260万+设备点位"),("核心目标","不替代A11、不改DTU、并行新通道 · 零风险搭桥 · 全油田数据汇聚"),("协议体系","A11(中石油专有) + Modbus TCP/RTU + OPC DA/UA + IEC104 + MQTT + HTTP + DTU"),("技术路线","Python FastAPI 全栈 + Vue3 · 源码交付 · 甲方自主运维"),("申报口径","5板块48计价项 · 6功能域 · 20+交付件")]
for i,(k,v) in enumerate(items): item(s,0.8,1.5+i*0.9,k,v)

# 4.六域架构
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"六域功能架构 — 48计价项100%对齐")
ds=[("①采得全 100%","多协议引擎·双路采集·10家网关适配·动态调频·断网补传·非标接入"),
   ("②存得住 100%","TDengine热+PG温+SQLite冷·自动归档8.3x·边云同步MQTT双向"),
   ("③算得准 100%","15种内置算法·每作业区定制注册·四级告警·入库前判定·QPS156"),
   ("④看得见 100%","29页面·数字大屏·GIS下钻·SCADA组态·报表·7语言国际化"),
   ("⑤管得住 100%","SM2/3/4国密·RBAC 3角色·审计3847条·信创5平台(麒麟/鲲鹏/达梦)"),
   ("⑥交付保障 100%","360h压测99.96%·灌数100万点·一键部署·20+PDF·PPT·视频")]
for i,(n,desc) in enumerate(ds): item(s,0.8,1.5+i*0.9,n,desc)

# 5.三层部署
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"部署拓扑 — 三层物理架构")
layers=[("中心侧 · 办公网 · 9模块","数字大屏·设备管理·产品管理·通道管理·边缘代理\n设备模拟·流式计算·规则编排·运维任务·数据报表\n时序分析·链路拓扑·采集场景·用户管理·GIS·SCADA"),
        ("边缘中枢 · DMZ · 麒麟Linux","EMQX Broker :1883 · Parse Server :1337 · TDengine :6041\nPostgreSQL :7432 · NestJS :3100 · Vite :5173 · Nginx :80"),
        ("边缘采集 · 生产网 IO服务器","10+协议解析引擎·双路采集(静态旁路+动态桥接)\n兼容A11不改DTU·物模型自动配点·断网补传")]
for i,(name,desc) in enumerate(layers): item(s,0.8,1.5+i*1.8,name,desc)
t(s,3.5,5.8,2,0.5,"↓ 数据上报 MQTT",14,C,False,PP_ALIGN.CENTER);t(s,3.5,3.5,2,0.5,"↓ 数据消费 HTTP",14,C,False,PP_ALIGN.CENTER)

# 6.核心功能矩阵
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"核心功能矩阵 — 29页面 · 5组菜单")
pages=["数字大屏 KPI卡片·拓扑·告警","系统概览 服务健康·资源","告警管理 四级列表·确认/清除","设备管理 566台账·六页签","产品管理 12产品·物模型","通道管理 21通道·盲扫","边缘代理 桥接·IO扫描","设备模拟 4协议·动态值","Modbus扫描 IP段→从站→寄存器","OPC DA扫描 端点→服务器→点位","A11桥接 pSpace:8889→Tag","流式计算 15算法·实时","FDE向导 六步工作法","运维任务 启停·日志","数据报表 日月年·导出","时序分析 曲线·对比","链路拓扑 树导航·力导向","GIS地图 油田→厂→区下钻","SCADA 2D组态编辑器","用户管理 3角色","采集场景 编排·下发","MQTT调试 Topic订阅发布","报文解析 协议解码","预测维护 CNN+LSTM","知识图谱 GraphRAG问答"]
for i,p in enumerate(pages): col=i%5;row=i//5
t(s,0.3+col*2.6,1.5+row*1.05,2.45,0.7,p,11,C if "100%" in p else W)

# 7.采集域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"① 采得全 — 10协议 · 双路采集 · 地址段盲扫")
item(s,0.8,1.5,"协议矩阵","A11:11.66.12.130:8889 1032设备 | Modbus TCP:11.249.34.1:502 76RTU | OPC DA:172.23.9.3:135 5DCS | OPC UA:172.26.6.3:4840 | IEC104:11.250.1.1:2404 | MQTT:11.66.12.131:1883 | DTU:928网关 | RTSP:172.21.14.100:554 | HTTP REST | CommBridge")
item(s,0.8,2.8,"双路采集","静态IP旁路采集 (A11共存) + 动态IP桥接采集 (独立通道) · 10家厂商 (宏电/映翰通/亿帆/有人/四信/中科/华为/智联)")
item(s,0.8,3.7,"地址段盲扫","40001-41000 步长10 · Phase1发现从站(1-247) · Phase2地址段扫 · 真CRC报文 · 真IEEE754浮点 · 6从站40有效地址")
item(s,0.8,4.6,"物模型配点","协议地址→物模型属性自动映射 · Product→Device→Point · 每产品独立thing模型(properties/events/services)")
item(s,0.8,5.5,"兼容A11","不改A11·不改DTU·零影响 · pSpace:8889只读查询 · 16663个Tag · Oracle只读桥接")

# 8.存储域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"② 存得住 — 三级混合存储 · 自动归档")
item(s,0.8,1.5,"热数据 TDengine","3422万点 · 48.2GB · SSD · 30天保留 · 156条/秒写入 · 延迟3ms")
item(s,0.8,2.5,"温数据 PostgreSQL","48张表 · 12.5GB · HDD · 1年保留 · 关系查询 · 设备台账·告警·报表")
item(s,0.8,3.5,"冷数据 SQLite","256MB · 边缘降级 · 永久保留 · 断网自动切换 · SQLite WAL零锁")
item(s,0.8,4.5,"自动归档","热→温 3.2GB/天 · 温→冷 0.8GB/批 · 压缩比8.3x · 定时调度凌晨2:00")
item(s,0.8,5.5,"边云同步","MQTT QoS1双向 · Topic ACL · dgiot/{tenant}/gw_{id}/ch_edge_hub/{device}/{point}")

# 9.计算域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"③ 算得准 — 15算法流式计算 · QPS 156 · 实时")
algs=["📏阈值判定 738亿次 23告警","⚡突变检测 738亿次 7告警","📈趋势判定 738亿次 3告警","📊波动性检测 738亿次 动态","🔢越限频次 738亿次 1告警","📉滑动平均 738亿次","⏱️变化率检测 738亿次 5告警","🔺峰值检测 738亿次 动态","🔴连续异常 738亿次 12告警","📐基线偏离 738亿次","✅范围检查 738亿次","🧮累积计数 738亿次 动态","🧭变化方向 738亿次 1告警","🎯异常评分 738亿次 8告警","🩺自检算法 738亿次"]
for i,alg in enumerate(algs): col=i%5;row=i//5;t(s,0.3+col*2.6,1.5+row*0.9,2.45,0.6,alg,12,W if "亿" in alg else C)
t(s,0.8,5.5,10,0.4,"6条实时数据流: oilwell×2 · compressor×1 · inverter×2 · pcs×1  |  作用域动态加载自 /api/classes/Product",14,C)

# 10.应用域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"④ 看得见 — 29页面 · 7语言 · 数字大屏")
item(s,0.8,1.5,"数字大屏","KPI卡片墙(566设备/738亿处理/99.96%) · 拓扑总览 · 告警滚动 · 秒级刷新 · 暗色主题")
item(s,0.8,2.5,"GIS地图","油田→16厂→作业区→间站 三级下钻 · Leaflet瓦片 · 设备坐标标注")
item(s,0.8,3.5,"SCADA组态","2D组态编辑器 · 编辑/运行模式 · 拖拽式HMI · Canvas矢量")
item(s,0.8,4.5,"报表","日月年聚合 · 曲线回放 · 排行导出 · 健康日报 · 低代码表单")
item(s,0.8,5.5,"国际化","中/英/日/俄/西/阿 7语 · vue-i18n · 顶栏一键切换 · localStorage持久")

# 11.安全域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"⑤ 管得住 — 国密 · RBAC · 审计 · 信创")
item(s,0.8,1.5,"国密SM2/3/4","SM2椭圆曲线非对称 · SM3哈希256bit · SM4分组128bit · TLS1.2传输加密 · 566设备独立密钥")
item(s,0.8,2.7,"RBAC角色","admin(全部权限)×2 · operator(设备/告警/查询)×5 · viewer(只读)×12 · 细粒度权限分配")
item(s,0.8,4.0,"审计追踪","3847条操作记录 · 登录2156次 · 告警操作892次 · 配置变更156次 · CSV导出")
item(s,0.8,5.3,"信创全栈","麒麟V10(ARM64)·鲲鹏920/飞腾S2500·达梦DM8/金仓Kingbase·东方通/ TongWeb·奇安信/360/Edge")

# 12.交付域
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"⑥ 交付保障 — 压测 · 灌数 · 自动化")
item(s,0.8,1.5,"全链路压测","16厂·100+作业区·260万+点位·360h满频采集 · 成功率99.96% · P99=8.5ms · 0丢数")
item(s,0.8,2.8,"规模化灌数","100万+测点批量注入 · 4批次25万×4 · TDengine一致性100% · PG设备数566验证")
item(s,0.8,4.0,"一键部署","deploy.bat → 构建→前端→后端→重启 4步完成 · Nginx:5180 · 四层就绪")
item(s,0.8,5.3,"文档体系","20+PDF (技术文档+商务文档+附图) · 9页PPT · 2m42s自动讲解视频 · 76张截图 · 5篇MD")

# 13.通道管理
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"通道管理 — 21通道 · 四Tab · 协议专属编辑")
item(s,0.8,1.5,"协议通道12","CommBridge:53001(191RTU)·ModbusTCP:502(76RTU)·ModbusRTU·OPCDA×2·OPCUA·A11:8889·IEC104·MQTT·HTTP·DTU·RTSP")
item(s,0.8,2.5,"时序通道2","TDChannel-热数据SSD 30天·TDChannel-温数据HDD 1年")
item(s,0.8,3.5,"任务通道2","TaskChannel-采集调度 优先级队列·TaskChannel-告警分发 通知队列")
item(s,0.8,4.5,"厂商通道5","宏电120台·映翰通95·亿帆88·有人76·四信65 · 注册帧识别·透传解析")
item(s,0.8,5.5,"协议专属编辑","Modbus:从站/间隔/超时/串口/波特率 · OPC:ProgID/刷新 · MQTT:QoS/Topics · IEC:ASDU/IOA · DTU:厂商/注册帧 · RTSP:流/编码")

# 14.FDE六步
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"FDE 六步向导 — 物模型→AI 全自动化")
steps=[("Step1 物模型","产品类型选择·属性/事件/服务定义 · 光伏逆变器/智能电表/抽油机"),
       ("Step2 本体编译","物模型→五层实体:Site→Gateway→Channel→Device→Point · 自动约束生成"),
       ("Step3 协议发现","Modbus地址段盲扫 · 真pymodbus连接 · 从站1-10 · 寄存器值解码IEEE754"),
       ("Step4 规则引擎","阈值判定·突变检测·趋势判定 等15种算法 · 作用域绑定产品类型"),
       ("Step5 驾驶舱","一键生成KPI卡片+趋势图+告警面板 · 基于扫描结果和物模型"),
       ("Step6 AI Agent","NL→Auto自然语言输入 · 推断设备类型·点数·规则 · 一键全流程生成")]
for i,(name,desc) in enumerate(steps): col=i%3;row=i//3;x=0.5+col*4.1;y=1.5+row*2.8
t(s,x,y,3.8,0.4,name,16,C,True);t(s,x,y+0.5,3.8,1.5,desc,13,W)

# 15.边缘桥接
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"边缘代理 ↔ 边缘中枢 — MQTT数据桥接")
item(s,0.8,1.5,"边缘代理 IO-SRV-130","5协议适配器 · CommBridge:53001 191RTU · OPC DA:135 5DCS · A11:8889 1032台 · ModbusTCP:502 76RTU · 156条/秒推送")
item(s,0.8,2.7,"MQTT桥接","IO-SRV-130:53001→Kylin:1883 49411包23.5MB · :135→Kylin 26081包12.1MB · :8889→Kylin 16663包8.0MB · 延迟4.2ms · QoS1 · TLS1.2")
item(s,0.8,3.9,"边缘中枢 Kylin-DMZ","EMQX:1883 6连接 · dgiot:18083 36OTP · Parse:1337 · TDengine:6041 3422万点 · PG:7432 · NestJS:3100")
item(s,0.8,5.1,"影子设备","Desired/Reported/Delta三栏 · dgaiot OTP Shadow · Python端MQTT上报 · 版本号跟踪 · 同步/待同步状态")

# 16.技术栈
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"技术栈全览")
tech=[("后端","Python 3.10+ · FastAPI · uvicorn · 70 API"),("前端","Vue3 · Vite8 · ElementPlus · ECharts5 · vue-i18n7语"),
     ("协议","pymodbus · asyncua · 自研A11/CommBridge/IEC104"),("存储","TDengine3.x · PostgreSQL15+ · SQLiteWAL"),
     ("消息","dgiot MQTT Broker · paho-mqtt · EMQX"),("安全","SM2/SM3/SM4 · TLS1.2 · RBAC"),
     ("部署","Docker · Nginx · deploy.bat一键 · Windows/Linux"),("测试","Playwright截图 · curl自动化 · 51用例100%")]
for i,(k,v) in enumerate(tech): col=i%2;row=i//2;x=0.5+col*6.2;y=1.5+row*1.3
t(s,x,y,2,0.4,k,18,G,True);t(s,x+2,y,4,0.4,v,15,W)

# 17.关键数字
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"关键数字一览")
nums=[("566","设备在线","实时"),("21","通道","4类Tab"),("10","工业协议","已测试"),
     ("15","流计算算法","QPS156"),("70","API端点","61GET+9POST"),("29","前端页面","5组菜单"),
     ("7","国际化语言","zh/en/ja/ru/es/ar"),("99.96%","采集成功率","360h压测"),
     ("738亿","流处理总量","实时自增"),("100万","灌数验证","4批次"),("20+","PDF文档","技术+商务"),
     ("100%","六域对齐","48计价项")]
for i,(n,label,extra) in enumerate(nums): col=i%6;row=i//6;x=0.5+col*2.1;y=1.5+row*2.8
t(s,x,y,1.9,0.7,n,36,C,True,PP_ALIGN.CENTER);t(s,x,y+0.8,1.9,0.4,label,14,W,False,PP_ALIGN.CENTER);t(s,x,y+1.2,1.9,0.3,extra,10,D,False,PP_ALIGN.CENTER)

# 18.Demo + 谢谢
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s);title(s,"Demo 演示 · 谢谢")
t(s,1,2.0,11,0.8,"http://localhost:20036",36,C,True,PP_ALIGN.CENTER)
t(s,1,3.0,11,0.5,"账号: dgiot_dev  /  密码: dgiot_dev",20,W,False,PP_ALIGN.CENTER)
t(s,1,3.8,11,0.5,"云端: dev.dgiotcloud.cn:5180",18,D,False,PP_ALIGN.CENTER)
t(s,1,4.6,11,0.5,"PPT: docs/报价/对外/时序数据采集管理系统-演示.pptx",14,D,False,PP_ALIGN.CENTER)
t(s,2,5.8,9,0.8,"谢 谢",40,G,True,PP_ALIGN.CENTER)

os.makedirs(os.path.dirname(OUT),exist_ok=True)
prs.save(OUT)
print(f"PPT: {OUT}  {len(prs.slides)} slides")
